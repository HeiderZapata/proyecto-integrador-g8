# -*- coding: utf-8 -*-
"""
Genera la PPT de defensa del PI Grupo 8 a partir del guion
(docs/guion_ppt_defensa.md). Cifras verificadas contra el repo.
Paleta y gramatica visual del curso SI7007. python-pptx 1.0.x.

Ejecutar:  python entrega_final/build_pptx.py
Salida:    entrega_final/PPT_Defensa_PI_Grupo8.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- Paleta (coherente con el tablero) ----------
STRUCT = RGBColor(0x2D, 0x2D, 0x2D)   # gris oscuro estructural
GREEN  = RGBColor(0x3B, 0x6D, 0x11)   # verde = oportunidad
RED    = RGBColor(0xA3, 0x2D, 0x2D)   # rojo = perdida/fuga
GRAY   = RGBColor(0xBB, 0xBB, 0xBB)   # gris = contexto
AMBER  = RGBColor(0xC8, 0x88, 0x1A)   # avisos "pendiente"
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
SOFT   = RGBColor(0x55, 0x55, 0x55)
LIGHTG = RGBColor(0xE9, 0xE9, 0xE9)
FONT   = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ---------- helpers ----------
def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_par(slide, l, t, w, h, text, size, color, bold=False, italic=False,
            align=PP_ALIGN.LEFT, anchor=None, font=FONT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = font
    return tb

def add_bullets(slide, l, t, w, h, items, size=16, color=STRUCT, gap=10):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        text, bold = (it, False) if isinstance(it, str) else it
        r = p.add_run(); r.text = "•  " + text
        f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = FONT
    return tb

def rect(slide, l, t, w, h, fill, line=None, rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is not None:
        shp.line.color.rgb = line; shp.line.width = Pt(1.25)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def box(slide, l, t, w, h, text, fill, text_color=WHITE, size=12, bold=True, line=None):
    shp = rect(slide, l, t, w, h, fill, line=line, rounded=True)
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = text_color; f.name = FONT
    return shp

def arrow(slide, cx, t, char="↓", size=16, color=SOFT):
    add_par(slide, cx - 0.3, t, 0.6, 0.3, char, size, color, align=PP_ALIGN.CENTER)

def footer(slide, presenter, time):
    add_par(slide, 0.6, 7.06, 8.0, 0.3, "PI Grupo 8 · Defensa 9-jun-2026", 9, GRAY)
    add_par(slide, 8.6, 7.06, 4.1, 0.3, f"{presenter} · {time}", 9, GRAY, align=PP_ALIGN.RIGHT)

def banner(slide, text, color):
    rect(slide, 0.6, 5.98, 12.13, 0.78, color, rounded=True)
    add_par(slide, 0.85, 6.04, 11.6, 0.66, text, 16, WHITE, bold=True, italic=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

def pending_callout(slide, text):
    rect(slide, 0.6, 6.0, 12.13, 0.72, AMBER, rounded=True)
    add_par(slide, 0.85, 6.05, 11.6, 0.62, text, 12.5, WHITE, bold=True,
            anchor=MSO_ANCHOR.MIDDLE)

def hbars(slide, x, y, w, h, items, maxval, caption=None):
    n = len(items)
    cap_h = 0.55 if caption else 0.0
    area_h = h - cap_h
    row_h = area_h / n
    bar_h = min(0.55, row_h * 0.55)
    label_w = w * 0.46
    bar_max_w = w - label_w - 1.05
    for i, (lab, val, col, vlab) in enumerate(items):
        cy = y + i * row_h + (row_h - bar_h) / 2.0
        add_par(slide, x, y + i * row_h, label_w, row_h, lab, 12.5, STRUCT,
                anchor=MSO_ANCHOR.MIDDLE)
        bw = max(0.04, bar_max_w * (val / maxval))
        rect(slide, x + label_w, cy, bw, bar_h, col)
        vcolor = STRUCT if col == GRAY else col
        add_par(slide, x + label_w + bw + 0.06, cy - 0.05, 1.0, bar_h + 0.1, vlab,
                12.5, vcolor, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if caption:
        add_par(slide, x, y + area_h + 0.04, w, cap_h, caption, 11, SOFT, italic=True)

# ---------- visuals ----------
def v_waffle(slide, x, y, w, h):
    step = 0.29; cell = 0.21
    grid = 10 * step
    sx = x + (w - grid) / 2.0
    sy = y + 0.02
    for idx in range(100):
        row, col = idx // 10, idx % 10
        c = GREEN if idx < 2 else GRAY
        rect(slide, sx + col * step, sy + row * step, cell, cell, c)
    add_par(slide, x, sy + grid + 0.05, w, 0.45, "2 de cada 100 compran  →  98 % se va sin comprar",
            13, RED, bold=True, align=PP_ALIGN.CENTER)

def v_funnel(slide, x, y, w, h):
    add_par(slide, x, y, w, 0.4, "Vistas  →  Carrito 3.86 %  →  Compra 2.27 %", 13, STRUCT,
            bold=True, align=PP_ALIGN.CENTER)
    barw = w - 0.6
    bx = x + 0.3; by = y + 1.0; bh = 1.05
    gw = barw * 0.5881; rw = barw * 0.4119
    rect(slide, bx, by, gw, bh, GREEN)
    rect(slide, bx + gw, by, rw, bh, RED)
    add_par(slide, bx, by, gw, bh, "Compran\n58.81 %", 13, WHITE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_par(slide, bx + gw, by, rw, bh, "Abandonan\n41.19 %", 13, WHITE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_par(slide, x, by + bh + 0.2, w, 0.5, "De cada carrito: 4 de cada 10 mueren ahi (903 k carritos · $250.4 M en juego)",
            11, SOFT, italic=True, align=PP_ALIGN.CENTER)

def v_recurrentes(slide, x, y, w, h):
    hbars(slide, x, y, w, h, [
        ("Recurrentes — 35.7 % de clientes · ticket $1,460", 73.7, GREEN, "73.7 %"),
        ("One-time — 64.3 % de clientes · ticket $289", 26.3, GRAY, "26.3 %"),
    ], 100, caption="% del revenue · 35.7 % de los compradores hacen el 73.7 % del dinero")

def v_models(slide, x, y, w, h):
    hbars(slide, x, y, w, h, [
        ("Dummy (azar) — trampa del accuracy", 0.056, GRAY, "0.056"),
        ("Logistica (baseline serio)", 0.097, GRAY, "0.097"),
        ("Familias RF / XGB / LGBM", 0.118, GRAY, "0.116–0.119"),
        ("LightGBM + Optuna (CV, 5 folds)", 0.1235, GREEN, "0.1235"),
    ], 0.14, caption="PR-AUC · test out-of-time 0.1172 ≈ CV 0.1235  →  sin fuga ni sobreajuste")

def v_lift(slide, x, y, w, h):
    hbars(slide, x, y, w, h * 0.72, [
        ("Base del test (prevalencia)", 5.6, GRAY, "5.6 %"),
        ("Grupo marcado por el modelo", 12.1, GREEN, "12.1 %"),
    ], 14, caption="Umbral F1-optimo 0.092 (no 0.5) → marca 15 % de sesiones · lift ~2.2×")
    add_par(slide, x, y + h * 0.80, w, 0.6,
            "Calibrado (Brier 0.0515): la probabilidad predicha ≈ frecuencia real de compra.",
            11.5, STRUCT, italic=True)

def v_importance(slide, x, y, w, h):
    add_par(slide, x, y, w, 0.4, "electronics_view_share — importancia", 13, STRUCT, bold=True)
    hbars(slide, x, y + 0.45, w, h * 0.5, [
        ("Built-in (split-count) — subestima", 0.042, GRAY, "0.042"),
        ("Permutacion (la que leemos)", 0.203, GREEN, "0.203"),
    ], 0.22)
    yy = y + 0.45 + h * 0.5 + 0.12
    add_par(slide, x, yy, w, 0.32, "Probabilidad de compra vs. precio de lo que miras (PDP):", 12, STRUCT, bold=True)
    ax_x = x + 1.0; ax_y0 = yy + 0.42; ax_y1 = yy + 1.45; ax_x1 = x + w - 0.3
    rect(slide, ax_x, ax_y0, 0.025, ax_y1 - ax_y0, GRAY)   # eje Y
    rect(slide, ax_x, ax_y1, ax_x1 - ax_x, 0.025, GRAY)    # eje X
    pts = 6
    for i in range(pts):
        px = ax_x + 0.18 + i * (ax_x1 - ax_x - 0.36) / (pts - 1)
        py = ax_y1 - 0.14 - i * (ax_y1 - ax_y0 - 0.22) / (pts - 1)
        rect(slide, px, py, 0.14, 0.14, GREEN)
    add_par(slide, x - 0.15, yy + 0.45, 1.05, 0.95, "P(compra) ↑", 10, SOFT, italic=True,
            anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    add_par(slide, ax_x, ax_y1 + 0.06, ax_x1 - ax_x, 0.3, "max_price_viewed  →", 10, SOFT,
            italic=True, align=PP_ALIGN.RIGHT)

def v_clusters(slide, x, y, w, h):
    hbars(slide, x, y, w, h, [
        ("C0 general bajo valor (49.8 % traf.)", 4.9, GRAY, "4.9 %"),
        ("C4 premium (13.8 % traf. · $1,024)", 5.8, GRAY, "5.8 %"),
        ("Base global (azar)", 5.97, GRAY, "5.97 %"),
        ("C3 gama media (29.3 % traf.) = OBJETIVO", 8.4, GREEN, "8.4 %"),
    ], 9.2, caption="% de conversion vs. base global 5.97 % · C3 la mayor. (C1 explorador 7 %, C2 anomalo)")

def v_data(slide, x, y, w, h):
    hbars(slide, x, y, w, h, [
        ("Eventos del clickstream", 109.95, GRAY, "109.95 M"),
        ("Sesiones (nuestra unidad)", 22.99, GREEN, "22.99 M"),
    ], 115, caption="1 fila = 1 sesion · base limpia de modelado 19.71 M (tasa de compra 5.97 %)")

def v_arch(slide, x, y, w, h):
    colw = 5.55; gap = w - 2 * colw
    lx = x + 0.1; rx = x + colw + gap - 0.1
    add_par(slide, lx, y, colw, 0.35, "REFERENCIA (produccion)", 13, SOFT, bold=True, align=PP_ALIGN.CENTER)
    add_par(slide, rx, y, colw, 0.35, "IMPLEMENTADA (Databricks Free)", 13, GREEN, bold=True, align=PP_ALIGN.CENTER)
    ref = ["Kafka / Kinesis", "Flink / Structured Streaming", "Delta sobre S3 / GCS",
           "Warehouse / Athena", "Serving en tiempo real"]
    imp = ["Volume (object storage)", "Auto Loader · Trigger.AvailableNow (Kappa)",
           "Bronze → Silver → Gold (Delta)", "Spark SQL", "MLflow + scoring batch",
           "Power BI (7 paginas, enlace publico)"]
    bh = 0.44; step = 0.50; top = y + 0.42
    for i, t in enumerate(ref):
        box(slide, lx, top + i * step, colw, bh, t, GRAY, text_color=STRUCT, size=11.5)
        if i < len(ref) - 1:
            arrow(slide, lx + colw / 2.0, top + i * step + bh - 0.02, size=12)
    box(slide, lx, top + 5 * step, colw, bh, "Gobernanza: Unity Catalog (permisos / accesos)",
        LIGHTG, text_color=STRUCT, size=10.5)
    for i, t in enumerate(imp):
        fill = GREEN if i in (1, 2, 4) else STRUCT
        box(slide, rx, top + i * step, colw, bh, t, fill, size=11.5)
        if i < len(imp) - 1:
            arrow(slide, rx + colw / 2.0, top + i * step + bh - 0.02, size=12)
    # frontera train/test como elemento marcado (rojo) bajo la columna implementada
    box(slide, rx + colw * 0.12, top + 5 * step + bh + 0.04, colw * 0.76, 0.32,
        "Gold: frontera TRAIN | TEST = filtro session_date", RED, size=9.5)

def v_cross(slide, x, y, w, h):
    bw = w - 1.0; bx = x + 0.5
    box(slide, bx, y + 0.05, bw, 0.95,
        "CLUSTERING  →  el DONDE:  bloque electronica C3 + C4 (43 % del trafico)",
        GREEN, size=14)
    arrow(slide, x + w / 2.0, y + 1.02, size=18)
    box(slide, bx, y + 1.35, bw, 0.95,
        "CLASIFICADOR  →  el A QUIEN:  ordena sesion a sesion dentro del bloque\n(Contrato 2: user_session · prob_calibrada · segmento)",
        STRUCT, size=13)
    arrow(slide, x + w / 2.0, y + 2.32, size=18)
    box(slide, bx + bw * 0.15, y + 2.62, bw * 0.7, 0.7,
        "= targeting fino: a quien mostrarle el incentivo", RED, size=14)

def v_ab(slide, x, y, w, h):
    bw = 2.55; bh = 0.95; top = y + 0.25
    xs = x + 0.1
    box(slide, xs, top, bw, bh, "Segmento\nC3", GREEN, size=13)
    arrow(slide, xs + bw + 0.18, top + bh / 2 - 0.15, char="→", size=18)
    box(slide, xs + bw + 0.36, top, bw, bh, "Aleatorizar\npor visitante", STRUCT, size=13)
    arrow(slide, xs + 2 * bw + 0.54, top + bh / 2 - 0.15, char="→", size=18)
    box(slide, xs + 2 * bw + 0.72, top, bw, bh * 0.46, "Control", GRAY, text_color=STRUCT, size=12)
    box(slide, xs + 2 * bw + 0.72, top + bh * 0.54, bw, bh * 0.46, "Tratamiento (incentivo)", STRUCT, size=11.5)
    arrow(slide, xs + 3 * bw + 0.9, top + bh / 2 - 0.15, char="→", size=18)
    box(slide, xs + 3 * bw + 1.08, top, bw + 0.15, bh, "Metrica:\nconv/visitante\n+ guardarrail margen", STRUCT, size=11)
    box(slide, x + 0.1, top + bh + 0.45, w - 0.2, 0.85,
        "Tamano:  MDE +0.5 pp (8.4 % → 8.9 %)  ⇒  ~49,600 por brazo  (~99,200 en total)  ·  alcanzable en pocos dias",
        GREEN, size=13)

def v_deploy(slide, x, y, w, h):
    bw = 2.75; bh = 1.0; top = y + 0.2; xs = x + 0.05
    seq = [("MLflow\n(modelo calibrado +\nsignature)", STRUCT),
           ("Scoring batch\n(pandas UDFs,\nsin OOM)", STRUCT),
           ("Contrato 2 en Delta\n19.71 M filas", GREEN),
           ("Alimenta:\nTablero + A/B", STRUCT)]
    for i, (t, c) in enumerate(seq):
        box(slide, xs + i * (bw + 0.18), top, bw, bh, t, c, size=11.5)
        if i < 3:
            arrow(slide, xs + (i + 1) * (bw + 0.18) - 0.09, top + bh / 2 - 0.15, char="→", size=16)
    box(slide, x + 0.05, top + bh + 0.4, w - 0.1, 0.8,
        "Verificado: al recargar reproduce PR-AUC 0.1172 · Brier 0.0515 · prob. media ≈ 0.060 (≈ tasa base)",
        GREEN, size=12.5)

def v_answers(slide, x, y, w, h):
    cw = (w - 0.5) / 2.0
    box(slide, x, y + 0.1, cw, 1.9,
        "DONDE\nEl carrito de electronica\n41 % de abandono · $187 M en juego", RED, size=14)
    box(slide, x + cw + 0.5, y + 0.1, cw, 1.9,
        "QUIEN\nC3 — electronica gama media\nmayor conversion (8.4 %)", GREEN, size=14)
    add_par(slide, x, y + 2.25, w, 0.7,
            "Palanca A: recuperar carritos de alta intencion   ·   Palanca B: retener al nucleo recurrente (73.7 % del revenue)",
            12.5, STRUCT, bold=True, align=PP_ALIGN.CENTER)

# ---------- slide scaffold ----------
def header(slide, num, title, message):
    add_par(slide, 0.62, 0.26, 1.2, 0.3, f"{num:02d}", 12, GRAY, bold=True)
    add_par(slide, 0.6, 0.5, 12.1, 0.8, title, 29, STRUCT, bold=True)
    rect(slide, 0.63, 1.3, 2.4, 0.05, GREEN)
    if message:
        add_par(slide, 0.6, 1.44, 12.1, 0.62, message, 17, GREEN, italic=True, bold=True)

def content_slide(num, title, message, presenter, time, bullets=None, visual=None,
                  layout="split", golden=None, pending=None, notes=""):
    s = prs.slides.add_slide(BLANK)
    _bg(s, WHITE)
    header(s, num, title, message)
    body_top = 2.25
    body_h = 3.45 if (golden or pending) else 4.55
    if layout == "split":
        if bullets:
            add_bullets(s, 0.6, body_top, 5.45, body_h, bullets, 15.5)
        if visual:
            visual(s, 6.35, body_top, 6.4, body_h)
    elif layout == "full":
        if visual:
            visual(s, 0.6, body_top, 12.13, body_h)
    elif layout == "bullets":
        if bullets:
            add_bullets(s, 0.6, body_top, 12.1, body_h, bullets, 18)
    if golden:
        banner(s, golden, GREEN)
    if pending:
        pending_callout(s, pending)
    footer(s, presenter, time)
    s.notes_slide.notes_text_frame.text = f"[{presenter} · {time}]\n{notes}"
    return s

def dark_slide(title, subtitle, lines, notes, presenter, time):
    s = prs.slides.add_slide(BLANK)
    _bg(s, STRUCT)
    add_par(s, 1.0, 2.2, 11.3, 1.6, title, 40, WHITE, bold=True)
    rect(s, 1.03, 3.7, 3.0, 0.06, GREEN)
    if subtitle:
        add_par(s, 1.0, 3.85, 11.3, 0.7, subtitle, 19, GRAY, italic=True)
    add_par(s, 1.0, 5.1, 11.3, 1.6, "\n".join(lines), 15, GRAY)
    s.notes_slide.notes_text_frame.text = f"[{presenter} · {time}]\n{notes}"
    return s

# =========================================================
#  SLIDES
# =========================================================

# 1 — Portada
dark_slide(
    "Optimizacion de la conversion en e-commerce",
    "Un sistema de decision sobre propension de compra · dataset REES46",
    ["Kelly  ·  Sara  ·  Heider  ·  Yeison",
     "Aprendizaje Automatico   ·   Grandes Datos   ·   Visualizacion",
     "EAFIT — Maestria en Ciencia de Datos  ·  9 de junio de 2026"],
    "Una frase de apertura: vamos a mostrar donde se fuga la conversion de la tienda y que hacer con eso. No leer la portada.",
    "Yeison", "0:15")

# 2 — El problema
content_slide(2, "El problema: la fuga de conversion",
    "La tienda pierde 98 de cada 100 visitas y 4 de cada 10 carritos.",
    "Yeison", "1:15",
    bullets=[("98 de cada 100 visitas no compran", True),
             "41 % de los carritos se abandonan",
             "La intencion existe; algo frena el cierre"],
    visual=v_waffle,
    golden="Pregunta de Oro — ¿Donde se concentra la fuga de conversion y que segmento de visitantes representa la mayor oportunidad de recuperarla?",
    notes=("De cada 100 personas que entran, 98 se van sin comprar; y de los que ponen algo en el "
           "carrito, 4 de cada 10 lo abandonan. No construimos un modelo para adivinar quien compra: "
           "construimos un SISTEMA DE DECISION para entender donde y por que se fuga la conversion, "
           "que visitantes la explican y donde intervenir. La prediccion es el instrumento; la decision "
           "de negocio es el producto. Plantar la Pregunta de Oro y prometer responderla al cierre."))

# 3 — Los datos
content_slide(3, "Los datos: escala Big Data, unidad = sesion",
    "Es un problema de Big Data real y modelamos por sesion, no por clic.",
    "Heider", "1:00",
    bullets=["REES46: clickstream multi-categoria, Oct–Nov 2019",
             ("109.95 M eventos  ·  ~14.5 GB crudos", True),
             "Gold = 22.99 M sesiones (1 fila = 1 user_session)",
             "Base limpia de modelado: 19.71 M · tasa 5.97 %"],
    visual=v_data,
    notes=("Trabajamos sobre un dataset real de ~110 millones de eventos y 14.5 GB. Lo importante: "
           "no modelamos eventos sueltos sino SESIONES, porque la pregunta de negocio es por visita, "
           "no por clic. Encuadre Big Data para el jurado de Grandes Datos."))

# 4 — Arquitectura
content_slide(4, "Arquitectura de datos: referencia vs. implementada",
    "Disenamos para produccion y lo implementamos acotado a Databricks Free.",
    "Heider", "1:45",
    layout="full", visual=v_arch,
    pending="PENDIENTE POR CONFIRMAR — Jurado de Grandes Datos: el perfil «BA/HPC» apunta a Edison Valencia; confirmar para dirigirle esta slide.",
    notes=("Plantilla 'escenario ideal vs. desarrollado' (lo premian): en produccion seria Kafka→Flink→S3; "
           "lo implementamos acotado en Databricks Free con un replay de streaming Kappa, que demuestra lo "
           "mismo sin broker externo. Particion POR TAMANO: la Gold no se particiona (1.33 GB; particionar "
           "daria micro-archivos de ~22 MB, el anti-patron) → OPTIMIZE + ZORDER + data-skipping. La frontera "
           "train/test se ve en el flujo (filtro session_date). Mencionar gobernanza con Unity Catalog en la "
           "referencia, sin implementarla."))

# 5 — EDA funnel
content_slide(5, "EDA · El embudo: donde se fuga",
    "La fuga grande no es solo 'no llegan al carrito': 4 de cada 10 carritos mueren ahi.",
    "Kelly", "0:50",
    bullets=["Vistas → carrito 3.86 % → compra 2.27 %",
             ("Abandono de carrito: 41.19 %", True),
             "903 k carritos · $250.4 M en juego"],
    visual=v_funnel,
    notes=("Declaracion + conector + razon: el abandono de carrito es 41 %, PORQUE la intencion existe "
           "pero algo frena el cierre. Ahi hay dinero recuperable sin traer mas trafico. Primera parte "
           "del 'donde' de la Pregunta de Oro."))

# 6 — EDA paradoja electronica
content_slide(6, "EDA · La paradoja de la electronica",
    "Casi todo el dinero en juego esta en una sola categoria.",
    "Kelly", "1:00",
    bullets=[("~$187 M de los $250 M en juego = electronica", True),
             "Mayor volumen y mayor pool de carritos abandonados",
             "Samsung + Apple = 68.6 % de los carritos de electronica"],
    visual=lambda s, x, y, w, h: hbars(s, x, y, w, h, [
        ("Resto (otras categorias + sin taxonomia)", 63, GRAY, "~$63 M"),
        ("Electronica", 187, RED, "$187 M"),
    ], 210, caption="Revenue en juego · electronica = 75 % de los $250.4 M en juego (incl. sin taxonomia)"),
    notes=("La oportunidad no esta repartida: se concentra en electronica, y dentro de ella en dos marcas. "
           "Eso enfoca donde mirar. CAVEAT para Q&A (no en slide): el $250 M es el total global incl. "
           "'Unknown'; las barras por-categoria excluyen 'Unknown' y suman ~$224 M. No invitar a sumarlas "
           "en vivo esperando $250 M."))

# 7 — EDA recurrentes
content_slide(7, "EDA · El nucleo recurrente",
    "Pocos compradores explican casi todo el revenue.",
    "Kelly", "0:50",
    bullets=[("35.7 % de compradores = 73.7 % del revenue", True),
             "Ticket recurrente $1,460 vs one-time $289 (5×)",
             "2a compra: mediana 2.9 dias · 84 % misma categoria"],
    visual=v_recurrentes,
    notes=("Aumento X debido a Y + pregunta: el revenue se concentra en los recurrentes PORQUE vuelven y "
           "gastan 5× mas. ¿Vale mas perseguir desconocidos o blindar a estos? Esa es la palanca B: retencion."))

# 8 — Estrategia de modelado
content_slide(8, "Estrategia de modelado",
    "La metrica y el split se eligen antes de competir modelos; el split, antes de las features.",
    "Sara", "1:30",
    bullets=["Target: ¿la sesion contiene purchase? (tasa 5.97 %)",
             ("Metrica: PR-AUC + Brier, NO accuracy (evento raro)", True),
             "Split temporal 'Opcion C': train Oct+Nov≤23 / test 24–30 nov",
             "Cuarentena 14–17 nov · features intra-sesion y pre-corte",
             ("El split se hace ANTES de construir features (anti-fuga)", True)],
    visual=lambda s, x, y, w, h: (
        rect(s, x + 0.3, y + 1.2, (w - 0.9) * 0.62, 0.7, GRAY),
        rect(s, x + 0.3 + (w - 0.9) * 0.62, y + 1.2, (w - 0.9) * 0.10, 0.7, RED),
        rect(s, x + 0.3 + (w - 0.9) * 0.72, y + 1.2, (w - 0.9) * 0.28, 0.7, GREEN),
        add_par(s, x + 0.3, y + 1.2, (w - 0.9) * 0.62, 0.7, "TRAIN\nOct + Nov ≤ 23", 12, STRUCT, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE),
        add_par(s, x + 0.3 + (w - 0.9) * 0.62, y + 0.55, (w - 0.9) * 0.12, 0.6, "cuarentena\n14–17", 9.5, RED, bold=True, align=PP_ALIGN.CENTER),
        add_par(s, x + 0.3 + (w - 0.9) * 0.72, y + 1.2, (w - 0.9) * 0.28, 0.7, "TEST\n24–30 nov", 12, WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE),
        add_par(s, x, y + 2.2, w, 0.5, "Un solo split, en la frontera · CV estratificada solo dentro del train", 11, SOFT, italic=True, align=PP_ALIGN.CENTER),
    ),
    notes=("Pre-emptar las TRES preguntas de Teran: TEMPORAL porque replica produccion (predecir el futuro); "
           "UNO SOLO en la frontera para no contaminar (la CV estratificada va dentro del train); ESE CORTE "
           "lo respalda la evidencia de deriva: PSI de precio ≈ 0 y conversion estable Oct↔Nov (notebook 03). "
           "Y lo clave: el split se hace ANTES de construir features, justo el error que les marcaron a pregrado."))

# 9 — Resultados: la cadena
content_slide(9, "Resultados: la cadena defendible",
    "Cada salto esta justificado y el test casi iguala a la validacion: sin fuga.",
    "Sara", "1:15",
    bullets=["Dummy 0.056 → Logistica 0.097 → Familias 0.116–0.119",
             ("LightGBM+Optuna: CV 0.1235 ± 0.0014 (5 folds)", True),
             "Test out-of-time 0.1172 (ex-Black-Friday 0.1115)",
             ("Test ≈ CV  ⇒  sin fuga ni sobreajuste", True)],
    visual=v_models,
    notes=("Frases-ancla de Teran: nunca decimos 'el mejor modelo' sin decir SEGUN QUE metrica, validacion y "
           "comparacion: aqui PR-AUC, CV de 5 folds, misma muestra. Como el test casi iguala a la CV, no hay "
           "fuga: 'si el test eligiera el modelo, dejaria de ser test', por eso el test lo tocamos una sola vez. "
           "Negocio: ordena a los compradores ~2× mejor que el azar."))

# 10 — Calibracion y umbral
content_slide(10, "Calibracion y umbral operativo",
    "El modelo estima probabilidades calibradas; el negocio decide el umbral, y no es 0.5.",
    "Sara", "1:00",
    bullets=[("Brier 0.0515 — probabilidades calibradas", True),
             "Umbral F1-optimo ≈ 0.092 (no 0.5)",
             "Marca 15 % de sesiones · precision 12.1 % · recall ~32 %",
             ("Lift ~2.2× sobre la base del test (5.6 %)", True)],
    visual=v_lift,
    notes=("Ancla de Teran: 'el modelo estima; el umbral decide'. 0.5 seria un error bajo desbalanceo; fijamos "
           "el umbral por costo/F1 en 0.092. Negocio: con el 15 % del esfuerzo se captura ~1/3 de las compras. "
           "CAVEAT (Q&A): la base 5.6 % es la prevalencia del test (nov 24–30); la base limpia global es 5.97 %. "
           "No mezclarlas."))

# 11 — Que predice la compra
content_slide(11, "Que predice la compra",
    "Lo que predice la compra es que tan caro/electronico es lo que miras, no cuanto navegas.",
    "Sara", "1:15",
    bullets=["Importancia por PERMUTACION (no la built-in del arbol)",
             ("max_price_viewed · electronics_view_share · avg_price_viewed", True),
             "La built-in subestima electronics_view_share (0.042 vs 0.203)",
             "Coherente con el EDA: se decide en ~2.2 min, navegan menos"],
    visual=v_importance,
    notes=("No repetimos el ranking: lo LEEMOS. La probabilidad SUBE con el precio y el peso de electronica de "
           "lo que ves, no con el numero de vistas. Por eso usamos permutacion y no la importancia nativa, que "
           "subestima a electronics_view_share por su baja cardinalidad. Cerrar con el insight contraintuitivo "
           "(2.2 min, navegan menos). OJO: no citar '2.2 min' cuando C3 este en pantalla (C3 perfila ~12 min)."))

# 12 — Segmentos (clustering)
content_slide(12, "Segmentos (clustering): aparece C3",
    "De 5 tipos de visitante, uno destaca por convertir y pesar a la vez: C3.",
    "Sara", "1:10",
    bullets=["k = 5 elegido por ACCIONABILIDAD (no por la metrica)",
             ("C3 'electronica gama media': 29.3 % traf. · conv 8.4 % (la mayor)", True),
             "C4 premium (13.8 %, $1,024, 5.8 %) · C0 bajo valor (49.8 %)",
             "C1 explorador (7 %) · C2 anomalo · DBSCAN: 1 masa + 3.7 % ruido"],
    visual=v_clusters,
    notes=("Ancla de Teran 'cluster != segmento': no llamamos 'segmento' a cualquier cluster: los nombramos y "
           "accionamos. Elegimos k=5 por ACCIONABILIDAD; el silhouette favoreceria k=2 (0.52) pero solo separa "
           "'electronica si/no', demasiado grueso para decidir. C3 = mayor conversion con volumen → objetivo a intervenir."))

# 13 — El cruce
content_slide(13, "El cruce: clasificador × clustering",
    "El segmento dice DONDE intervenir; el modelo dice A QUIEN.",
    "Yeison", "1:05",
    layout="full", visual=v_cross,
    notes=("El corazon del proyecto: el clustering dice DONDE esta la oportunidad (electronica) y el modelo de "
           "propension, dentro de ahi, A QUIEN priorizar sesion a sesion. Ninguno solo basta; el cruce es el insight."))

# 14 — A/B
content_slide(14, "El experimento (A/B)",
    "Con datos observacionales no prometemos uplift causal: entregamos el DISENO.",
    "Yeison", "1:15",
    layout="full", visual=v_ab,
    notes=("Cierra el alcance (responde al revisor): el dato es observacional, no hay variable de tratamiento, "
           "asi que estimar uplift seria afirmar causalidad no verificable. Lo honesto es DISENAR el A/B que si "
           "lo mediria: sobre C3, aleatorizando por visitante, con guardarrail de margen. Con un MDE de medio "
           "punto, ~49,600 visitantes por brazo, se corre en dias."))

# 15 — Despliegue
content_slide(15, "Despliegue tecnologico",
    "El modelo no quedo en un notebook: se persiste, alimenta el tablero y el A/B, y esta verificado.",
    "Heider", "1:00",
    layout="full", visual=v_deploy,
    notes=("Cerramos el ciclo de vida del dato: el modelo se loguea en MLflow y se aplica a las 19.7 M sesiones; "
           "el resultado (el Contrato 2) vive en Delta y alimenta el tablero y el experimento. Y lo validamos: al "
           "recargarlo, las probabilidades reproducen PR-AUC y Brier, y su media coincide con la tasa base. "
           "Cubre 'persistencia de modelos' de Grandes Datos."))

# 16 — Conclusiones
content_slide(16, "Conclusiones: respondemos la Pregunta de Oro",
    "La fuga se concentra en el carrito de electronica; C3 es la mayor oportunidad.",
    "Yeison", "1:20",
    layout="full", visual=v_answers,
    golden="Pregunta de Oro — respondida: DONDE = carrito de electronica · QUIEN = C3 (electronica gama media)",
    notes=("Cerrar el arco: abrimos preguntando donde se fuga la conversion y que segmento la explica. Respuesta: "
           "se fuga en el CARRITO DE ELECTRONICA y el segmento de mayor oportunidad es C3. El modelo prioriza a "
           "quien, el experimento mediria cuanto. Aporte por materia: ML (propension + clustering calibrados), "
           "Grandes Datos (Medallion/Kappa + scoring + Contrato 2), Visualizacion (tablero ejecutivo desplegado). "
           "La prediccion fue el instrumento; la decision de negocio es el producto."))

# 17 — Demo
content_slide(17, "Demo del tablero en vivo",
    "Lo que contamos se sostiene con datos, en vivo, desde el enlace publico.",
    "Kelly", "2:00",
    bullets=["Power BI Service · enlace publico · 7 paginas",
             "Recorrido: Problema → Detalle Electronics → Cliente → Cierre",
             ("Filtrar EN VIVO (Top N marcas, categoria)", True),
             "Responder la Pregunta de Oro desde el dashboard"],
    visual=lambda s, x, y, w, h: box(s, x + 1.2, y + 0.6, w - 2.4, 2.2,
        "TABLERO EN VIVO\n(abrir el enlace publico y filtrar)\n\nTener captura de respaldo por si falla la red",
        STRUCT, size=14),
    pending="PENDIENTE POR CONFIRMAR — ¿El tablero ya muestra el clustering (C3) en 'Analisis de cliente'? Si no, enrutar el demo por 'Detalle Electronics'.",
    notes=("Esto esta desplegado y es publico. Filtro electronica... aqui esta el pool de carritos en juego; "
           "filtro por marca... Samsung y Apple. La historia de las slides es exactamente lo que ven aqui. "
           "Definir publico (negocio/direccion) y objetivo (convencer de donde intervenir). "
           "Ensayar el filtrado antes; tener captura de respaldo."))

# 18 — Cierre
dark_slide(
    "Gracias",
    "¿Donde se concentra la fuga y que segmento es la mayor oportunidad? — lista para sus preguntas",
    ["La prediccion es el instrumento; la decision de negocio es el producto.",
     "PI Grupo 8  ·  Kelly · Sara · Heider · Yeison"],
    "Invitar preguntas; al responder, volver al tablero o a la slide concreta que respalda. No responder de memoria.",
    "Equipo", "Q&A")

# =========================================================
import os
out = os.path.join(os.path.dirname(__file__), "PPT_Defensa_PI_Grupo8.pptx")
prs.save(out)
print(f"OK · {len(prs.slides._sldIdLst)} slides · {out}")
