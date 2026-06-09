# -*- coding: utf-8 -*-
"""Vuelca el contenido del .pptx (texto por shape + notas) a un .txt auditable."""
import os
from pptx import Presentation
from pptx.util import Emu

here = os.path.dirname(__file__)
prs = Presentation(os.path.join(here, "PPT_Defensa_PI_Grupo8.pptx"))
lines = []
for i, slide in enumerate(prs.slides, 1):
    lines.append(f"\n========== SLIDE {i} ==========")
    for sh in slide.shapes:
        kind = sh.shape_type
        try:
            l = round(Emu(sh.left).inches, 2); t = round(Emu(sh.top).inches, 2)
            w = round(Emu(sh.width).inches, 2); h = round(Emu(sh.height).inches, 2)
            geo = f"[{l},{t} {w}x{h}]"
        except Exception:
            geo = "[-]"
        if sh.has_text_frame and sh.text_frame.text.strip():
            txt = sh.text_frame.text.replace("\n", " ⏎ ")
            lines.append(f"  TXT {geo}: {txt}")
    notes = ""
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
    lines.append(f"  >> NOTAS: {notes}")
out = os.path.join(here, "_pptx_dump.txt")
with open(out, "w", encoding="utf-8") as f:
    f.write("\n".join(lines))
print(f"dump OK -> {out} ({len(prs.slides._sldIdLst)} slides)")
